"""Stage 1 of the KB ingestion pipeline: extract text and metadata from a source binary.

Given an input file (PDF, DOCX, EPUB, PPTX, MD, TXT), this script:
  1. Detects the format and an appropriate extractor.
  2. Picks a slug (from filename, sanitized) or uses --slug if provided.
  3. Creates raw/<slug>/ with the original binary copied in as original.<ext>.
  4. Extracts the text into raw/<slug>/source.md with locator markers.
  5. Writes raw/<slug>/source.json with metadata.
  6. Initializes raw/<slug>/log.md with a "cataloged" entry.

This script does NOT do any LLM work. It's deterministic, free of LLM tokens,
and produces the canonical inputs the librarian agent reads in later stages.

Run via uv:

    uv run python tooling/scripts/extract.py <path-to-source>

Optional flags:
    --slug <name>              override the auto-generated slug
    --kb-root <path>           KB root directory (defaults to cwd)
    --extractor <name>         force a specific extractor (pypdf, pdfplumber,
                               opendataloader, ebooklib, docx, pptx, copy)
"""

from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import json
import re
import shutil
import sys
from pathlib import Path
from typing import Callable

EXTRACTOR_VERSION = "0.1.0"


# -------- slug + format detection --------

FORMAT_BY_EXT = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".epub": "epub",
    ".pptx": "pptx",
    ".md": "md",
    ".markdown": "md",
    ".txt": "txt",
}


def detect_format(path: Path) -> str:
    suffix = path.suffix.lower()
    fmt = FORMAT_BY_EXT.get(suffix)
    if not fmt:
        raise ValueError(f"unsupported format: {suffix} (file: {path.name})")
    return fmt


def slugify(name: str) -> str:
    """Turn a filename like 'Cheng et al. (2026) - Sycophantic AI.pdf' into
    'cheng-et-al-2026-sycophantic-ai'. Non-rigorous; the librarian may
    rename to a canonical author-keyword-year stem at integration time."""
    stem = Path(name).stem
    # Remove parentheses contents but keep year if present
    stem = re.sub(r"\((\d{4})\)", r"\1", stem)
    stem = re.sub(r"\([^)]*\)", "", stem)
    # Replace separators
    stem = re.sub(r"[\s_]+", "-", stem)
    # Lowercase, collapse ascii punctuation
    stem = stem.lower()
    stem = re.sub(r"[^a-z0-9-]+", "-", stem)
    stem = re.sub(r"-+", "-", stem).strip("-")
    return stem


def heuristic_source_type(filename: str, format: str, first_chars: str) -> str:
    """Best-effort guess at source_type. Librarian can override during triage."""
    fn = filename.lower()
    text = first_chars.lower()
    if format in {"epub"}:
        return "book"
    if format in {"pptx"}:
        return "transcript"  # decks usually transcribe a talk
    if any(k in fn for k in ["transcript", "podcast", "interview", "video"]):
        return "transcript"
    if any(k in fn for k in ["report", "barometer", "survey"]) or "executive summary" in text:
        return "report"
    # Heuristic for academic papers vs articles
    if format == "pdf":
        # Strong paper signals: classic IMRaD section headings (opendataloader-pdf
        # extracts these as markdown headings), DOI/arXiv prefixes, abstract block
        paper_signals = [
            "abstract", "doi:", "arxiv:",
            "## introduction", "### introduction",
            "## method", "### method", "## methods", "### methods",
            "## discussion", "### discussion",
            "## conclusion", "### conclusion",
            "introduction\n", "references\n",
        ]
        if any(k in text for k in paper_signals):
            return "paper"
        return "article"
    if format in {"md", "txt"}:
        return "article"
    return "unknown"


# -------- extractors --------

def extract_pdf_pypdf(path: Path) -> tuple[str, dict]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        pages.append(f"\n\n[p.{i}]\n\n{text.strip()}")
    body = "\n".join(pages).strip()
    meta = {"pages": len(reader.pages)}
    return body, meta


def extract_pdf_pdfplumber(path: Path) -> tuple[str, dict]:
    import pdfplumber

    pages = []
    page_count = 0
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_count = i
            text = page.extract_text() or ""
            pages.append(f"\n\n[p.{i}]\n\n{text.strip()}")
    body = "\n".join(pages).strip()
    return body, {"pages": page_count}


def extract_pdf_opendataloader(path: Path) -> tuple[str, dict]:
    """Primary PDF extractor when Java is available. Preserves bounding
    boxes, page numbers, heading hierarchy, element types via the
    Java-backed structural extractor."""
    try:
        from opendataloader_pdf import convert  # type: ignore
    except Exception as e:  # pragma: no cover - optional dep
        raise RuntimeError(f"opendataloader-pdf not available: {e}")

    import tempfile
    with tempfile.TemporaryDirectory() as td:
        convert(
            str(path),
            output_dir=td,
            format="markdown",
            quiet=True,
            markdown_page_separator="\n\n[p.%page-number%]\n\n",
        )
        md_files = list(Path(td).glob("*.md"))
        if not md_files:
            raise RuntimeError("opendataloader-pdf produced no markdown output")
        body = md_files[0].read_text(encoding="utf-8")

    # Count pages from the locator markers we requested
    page_markers = re.findall(r"\[p\.(\d+)\]", body)
    pages = max((int(p) for p in page_markers), default=0)
    return body, {"pages": pages}


def extract_docx(path: Path) -> tuple[str, dict]:
    from docx import Document

    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name or "").lower()
        if style.startswith("heading"):
            try:
                level = int(re.search(r"\d+", style).group())
            except Exception:
                level = 2
            parts.append("\n" + ("#" * max(1, min(6, level))) + " " + text + "\n")
        else:
            parts.append(text)
    body = "\n\n".join(parts).strip()
    return body, {"paragraphs": len(doc.paragraphs)}


def extract_epub(path: Path) -> tuple[str, dict]:
    from ebooklib import epub, ITEM_DOCUMENT
    from html import unescape

    book = epub.read_epub(str(path))
    chapters: list[str] = []
    chap_no = 0
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        chap_no += 1
        raw = item.get_content().decode("utf-8", errors="replace")
        # Strip HTML tags lightly; this is a coarse extraction. The
        # librarian only needs readable text + chapter locators.
        text = re.sub(r"<[^>]+>", " ", raw)
        text = unescape(text)
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            continue
        title = item.get_name() or f"chapter-{chap_no}"
        chapters.append(f"\n\n[ch.{chap_no} — {title}]\n\n{text}")
    body = "\n".join(chapters).strip()
    return body, {"chapters": chap_no}


def extract_pptx(path: Path) -> tuple[str, dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        bits = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        bits.append(text)
        slides.append(f"\n\n[slide {i}]\n\n" + "\n".join(bits))
    body = "\n".join(slides).strip()
    return body, {"slides": len(prs.slides)}


def extract_markdown(path: Path) -> tuple[str, dict]:
    body = path.read_text(encoding="utf-8", errors="replace")
    lines = body.count("\n") + 1
    return body, {"lines": lines}


# -------- driver --------

EXTRACTOR_CHAIN: dict[str, list[tuple[str, Callable[[Path], tuple[str, dict]]]]] = {
    "pdf": [
        ("opendataloader-pdf", extract_pdf_opendataloader),
        ("pypdf", extract_pdf_pypdf),
        ("pdfplumber", extract_pdf_pdfplumber),
    ],
    "docx": [("python-docx", extract_docx)],
    "epub": [("ebooklib", extract_epub)],
    "pptx": [("python-pptx", extract_pptx)],
    "md": [("copy", extract_markdown)],
    "txt": [("copy", extract_markdown)],
}


def run_extractor_chain(format: str, path: Path, force: str | None) -> tuple[str, str, dict]:
    chain = EXTRACTOR_CHAIN[format]
    last_err: Exception | None = None
    for name, fn in chain:
        if force and force != name:
            continue
        try:
            body, meta = fn(path)
            if body and len(body) > 200:
                return name, body, meta
            # extractor produced something too small — try next
            last_err = RuntimeError(f"{name} returned only {len(body)} chars")
        except Exception as e:
            last_err = e
    raise RuntimeError(f"all extractors failed for {path.name}: {last_err}")


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def reextract_existing(slug: str, kb_root: Path, forced_extractor: str | None) -> int:
    """Re-run extraction on an existing raw/<slug>/. Updates source.md and the
    extractor/word_count/pages fields in source.json. Preserves slug, status,
    decisions, ingested_at, and any backfill markers."""
    out_dir = (kb_root / "raw" / slug).resolve()
    sj_path = out_dir / "source.json"
    if not sj_path.is_file():
        print(f"error: {sj_path} not found (no existing extraction to upgrade)", file=sys.stderr)
        return 4
    meta = json.loads(sj_path.read_text(encoding="utf-8"))
    fmt = meta.get("format")
    if not fmt:
        print(f"error: source.json has no format field", file=sys.stderr)
        return 5
    original = next(out_dir.glob(f"original.*"), None)
    if original is None:
        print(f"error: no original.* binary found in {out_dir}", file=sys.stderr)
        return 6

    old_extractor = meta.get("extractor", "unknown")
    extractor_used, body, extra_meta = run_extractor_chain(fmt, original, forced_extractor)
    (out_dir / "source.md").write_text(body, encoding="utf-8")

    word_count = len(re.findall(r"\b\w+\b", body))
    now = dt.datetime.now(dt.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

    # Update only the extraction-related fields; preserve everything else.
    meta["extractor"] = extractor_used
    meta["extractor_script_version"] = EXTRACTOR_VERSION
    meta["word_count"] = word_count
    for k, v in extra_meta.items():
        meta[k] = v
    meta.setdefault("decisions", []).append({
        "stage": "reextract",
        "outcome": "extractor_upgraded",
        "by": "tooling/scripts/extract.py --reextract",
        "at": now,
        "reason": f"Re-extracted with {extractor_used} (was {old_extractor}). source.md replaced; lifecycle state preserved.",
    })
    sj_path.write_text(json.dumps(meta, indent=2), encoding="utf-8")

    log_path = out_dir / "log.md"
    note = (
        f"\n## {now} — re-extracted\n"
        f"- Extractor: {extractor_used} (was {old_extractor})\n"
        f"- Word count: {word_count}\n"
    )
    with log_path.open("a", encoding="utf-8") as f:
        f.write(note)

    print(f"re-extracted: {out_dir}")
    print(f"  extractor: {old_extractor} -> {extractor_used}")
    print(f"  word_count: {word_count}")
    return 0


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract a source for the Modrn Mind KB")
    parser.add_argument("path", nargs="?", type=Path, help="path to the source file (omit when using --reextract)")
    parser.add_argument("--slug", help="override auto-generated slug, OR (with --reextract) the slug to re-extract")
    parser.add_argument("--kb-root", type=Path, default=Path.cwd(), help="KB root (default: cwd)")
    parser.add_argument("--extractor", help="force a specific extractor")
    parser.add_argument("--force", action="store_true", help="overwrite existing raw/<slug>/")
    parser.add_argument("--reextract", action="store_true", help="re-run extraction on an existing raw/<slug>/; preserves lifecycle state")
    args = parser.parse_args()

    if args.reextract:
        if not args.slug:
            print("error: --reextract requires --slug", file=sys.stderr)
            return 2
        return reextract_existing(args.slug, args.kb_root.resolve(), args.extractor)

    if args.path is None:
        print("error: path argument required (unless --reextract)", file=sys.stderr)
        return 2

    src: Path = args.path.resolve()
    if not src.is_file():
        print(f"error: {src} is not a file", file=sys.stderr)
        return 2

    fmt = detect_format(src)
    slug = args.slug or slugify(src.name)

    out_dir = (args.kb_root / "raw" / slug).resolve()
    if out_dir.exists() and not args.force:
        print(f"error: {out_dir} already exists (use --force to overwrite, or --reextract to upgrade in place)", file=sys.stderr)
        return 3
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1. Copy original
    original = out_dir / f"original{src.suffix.lower()}"
    shutil.copy2(src, original)

    # 2. Run extraction
    extractor_used, body, extra_meta = run_extractor_chain(fmt, original, args.extractor)

    # 3. Write source.md
    (out_dir / "source.md").write_text(body, encoding="utf-8")

    # 4. Write source.json
    word_count = len(re.findall(r"\b\w+\b", body))
    first_chars = body[:2000]
    src_type = heuristic_source_type(src.name, fmt, first_chars)
    now = dt.datetime.now(dt.timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    meta = {
        "slug": slug,
        "original_filename": src.name,
        "format": fmt,
        "source_type": src_type,
        "word_count": word_count,
        "hash": file_hash(original),
        "extractor": extractor_used,
        "extractor_script_version": EXTRACTOR_VERSION,
        "status": "cataloged",
        "ingested_at": now,
        "decisions": [],
        **extra_meta,
    }
    (out_dir / "source.json").write_text(json.dumps(meta, indent=2), encoding="utf-8")

    # 5. Initialize log.md
    log = (
        f"# Log: {slug}\n\n"
        f"## {now} — cataloged\n"
        f"- Extractor: {extractor_used}\n"
        f"- Format: {fmt} ({extra_meta})\n"
        f"- Word count: {word_count}\n"
        f"- Heuristic source_type: {src_type}\n"
    )
    (out_dir / "log.md").write_text(log, encoding="utf-8")

    print(f"cataloged: {out_dir}")
    print(f"  extractor: {extractor_used}")
    print(f"  word_count: {word_count}")
    print(f"  source_type: {src_type}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
